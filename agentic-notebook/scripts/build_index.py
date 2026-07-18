#!/usr/bin/env python3
"""
build_index.py — Turn a folder of files into a NotebookLM-style source index.

Creates <folder>/.notebook/ containing:
  sources.json   registry of every source (stable IDs S1, S2, ...)
  chunks.jsonl   one JSON object per retrievable chunk (IDs like S1#004)
  meta.json      index metadata
  studio/        (empty) generated artifacts go here
  notes/         (empty) saved notes go here

Usage:
  python build_index.py <folder>                      build or refresh the index
  python build_index.py <folder> --status             show index status
  python build_index.py <folder> --from-text <relpath> <textfile>
        Register text the agent extracted itself (for scanned PDFs, images,
        audio transcripts, or anything the script could not parse).
        Inside <textfile>, lines like "[[page 12]]" mark page boundaries.

Pure standard library at its core. PDF/DOCX/PPTX extraction uses optional
libraries when present (pymupdf, pdfplumber, pypdf, python-docx, python-pptx)
and falls back to stdlib zip/XML parsing or the `pdftotext` CLI. Files that
cannot be parsed are registered with status "needs_extraction" so the agent
can read them natively and feed the text back with --from-text.
"""

import os
import sys
from pathlib import Path

# --- Environment bootstrap: re-exec inside <skill>/.venv if present --------
# A user/agent may create a virtual environment at <skill>/.venv to hold
# optional dependencies (pymupdf, easyocr, ...) on machines where the system
# Python can't install packages. Detect it and hand execution over.
_venv = Path(__file__).resolve().parent.parent / ".venv"
_venv_py = _venv / ("Scripts/python.exe" if os.name == "nt" else "bin/python")
try:
    # sys.prefix equals the venv dir when running inside it — comparing
    # executables fails on Linux, where venv pythons are symlinks.
    _needs_venv = (_venv_py.exists()
                   and Path(sys.prefix).resolve() != _venv.resolve())
except OSError:
    _needs_venv = False
if _needs_venv:
    import subprocess as _sp
    raise SystemExit(_sp.call([str(_venv_py)] + sys.argv))
# ---------------------------------------------------------------------------
import argparse
import csv
import hashlib
import html
import io
import json
import os
import re
import subprocess
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
try:
    from ocr_pdf import run_ocr, pick_engines  # noqa: E402
except Exception:  # keep build_index usable even if ocr_pdf.py is absent
    run_ocr = None
    pick_engines = None

NOTEBOOK_DIR = ".notebook"
INDEX_VERSION = 3  # bump when extraction/chunking changes; forces re-extraction
SCAN_WORDS_PER_PAGE = 15   # a PDF averaging fewer words/page than this is treated as scanned
CHUNK_TARGET_WORDS = 230
CHUNK_MAX_WORDS = 340
CHUNK_MIN_WORDS = 60

TEXT_EXTS = {".txt", ".md", ".markdown", ".rst", ".log", ".tex", ".org",
             ".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs", ".rb",
             ".sh", ".sql", ".yaml", ".yml", ".toml", ".ini", ".srt", ".vtt"}
HTML_EXTS = {".html", ".htm", ".xhtml"}
CSV_EXTS = {".csv", ".tsv"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}
EPUB_EXTS = {".epub"}
JSON_EXTS = {".json", ".jsonl", ".ndjson"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".wma"}
VIDEO_EXTS = {".mp4", ".mov", ".mkv", ".webm", ".avi"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}
SKIP_EXTS = {".zip", ".tar", ".gz", ".7z", ".rar", ".exe", ".dll", ".so",
             ".bin", ".dat", ".db", ".sqlite", ".pkl", ".pt", ".onnx",
             ".doc", ".xls", ".xlsx", ".ppt"}  # legacy office: agent extracts
MAX_FILE_BYTES = 250 * 1024 * 1024


# ---------------------------------------------------------------- utilities

_PIPE_CLOSED = False


def say(*args, **kw):
    """print() that survives a closed stdout (e.g. piped into `head`).
    Progress output is disposable; the index write is not — never let a
    BrokenPipeError abort indexing before save_all runs."""
    global _PIPE_CLOSED
    if _PIPE_CLOSED:
        return
    try:
        print(*args, **kw)
    except BrokenPipeError:
        _PIPE_CLOSED = True


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(1 << 20), b""):
            h.update(block)
    return h.hexdigest()


def norm_ws(text: str) -> str:
    text = text.replace("\u00ad", "")           # soft hyphens
    # Rejoin words hyphenated across line breaks ("vol-\nume" -> "volume"),
    # which plain PDF extraction leaves split and which breaks search tokens.
    text = re.sub(r"([A-Za-z\u00c0-\u024f])-\n([a-z\u00e0-\u024f])", r"\1\2", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def title_from_path(rel: str) -> str:
    stem = Path(rel).stem
    stem = re.sub(r"[_\-]+", " ", stem).strip()
    return stem or rel


def strip_html(raw: str) -> str:
    """Minimal, dependency-free HTML → text."""
    raw = re.sub(r"(?is)<(script|style|nav|header|footer)[^>]*>.*?</\1>", " ", raw)
    raw = re.sub(r"(?i)<br\s*/?>", "\n", raw)
    raw = re.sub(r"(?i)</(p|div|h[1-6]|li|tr|section|article|blockquote)>", "\n\n", raw)
    raw = re.sub(r"(?s)<[^>]+>", " ", raw)
    return norm_ws(html.unescape(raw))


# ---------------------------------------------------------- extractors
# Every extractor returns a list of (unit_label, text) pairs, where
# unit_label is something like "p.3", "slide 2", or "" for unpaged text —
# or raises / returns None when it cannot handle the file.

def extract_plain(path: Path):
    text = path.read_text(encoding="utf-8", errors="replace")
    return [("", norm_ws(text))]


def extract_html_file(path: Path):
    return [("", strip_html(path.read_text(encoding="utf-8", errors="replace")))]


def extract_csv_file(path: Path):
    delim = "\t" if path.suffix.lower() == ".tsv" else ","
    rows = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for i, row in enumerate(csv.reader(f, delimiter=delim)):
            rows.append(" | ".join(cell.strip() for cell in row))
            if i > 20000:
                rows.append("... (truncated: file has more rows)")
                break
    return [("", norm_ws("\n".join(rows)))]


def extract_json_file(path: Path):
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        obj = json.loads(raw)
        raw = json.dumps(obj, indent=1, ensure_ascii=False)
    except Exception:
        pass
    if len(raw) > 2_000_000:
        raw = raw[:2_000_000] + "\n... (truncated)"
    return [("", norm_ws(raw))]


def extract_pdf(path: Path):
    # 1. PyMuPDF
    try:
        import fitz  # type: ignore
        doc = fitz.open(str(path))
        pages = [(f"p.{i + 1}", norm_ws(page.get_text("text"))) for i, page in enumerate(doc)]
        doc.close()
        if any(t for _, t in pages):
            return pages, "pymupdf"
    except Exception:
        pass
    # 2. pdfplumber
    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(str(path)) as pdf:
            pages = [(f"p.{i + 1}", norm_ws(pg.extract_text() or "")) for i, pg in enumerate(pdf.pages)]
        if any(t for _, t in pages):
            return pages, "pdfplumber"
    except Exception:
        pass
    # 3. pypdf
    try:
        import pypdf  # type: ignore
        reader = pypdf.PdfReader(str(path))
        pages = [(f"p.{i + 1}", norm_ws(pg.extract_text() or "")) for i, pg in enumerate(reader.pages)]
        if any(t for _, t in pages):
            return pages, "pypdf"
    except Exception:
        pass
    # 4. pdftotext CLI
    try:
        out = subprocess.run(["pdftotext", "-enc", "UTF-8", str(path), "-"],
                             capture_output=True, timeout=300)
        if out.returncode == 0 and out.stdout.strip():
            raw = out.stdout.decode("utf-8", errors="replace")
            pages = [(f"p.{i + 1}", norm_ws(t)) for i, t in enumerate(raw.split("\f"))]
            pages = [p for p in pages if p[1]]
            if pages:
                return pages, "pdftotext"
    except Exception:
        pass
    return None, None


def extract_docx(path: Path):
    try:
        import docx  # type: ignore
        d = docx.Document(str(path))
        parts = [p.text for p in d.paragraphs]
        for table in d.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells))
        return [("", norm_ws("\n\n".join(x for x in parts if x.strip())))], "python-docx"
    except Exception:
        pass
    try:  # stdlib fallback: unzip the XML
        with zipfile.ZipFile(path) as z:
            raw = z.read("word/document.xml").decode("utf-8", errors="replace")
        raw = re.sub(r"</w:p>", "\n\n", raw)
        raw = re.sub(r"<[^>]+>", "", raw)
        return [("", norm_ws(html.unescape(raw)))], "zip-xml"
    except Exception:
        return None, None


def extract_pptx(path: Path):
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                nt = slide.notes_slide.notes_text_frame.text
                if nt.strip():
                    texts.append("[speaker notes] " + nt)
            slides.append((f"slide {i}", norm_ws("\n".join(t for t in texts if t.strip()))))
        return slides, "python-pptx"
    except Exception:
        pass
    try:  # stdlib fallback
        slides = []
        with zipfile.ZipFile(path) as z:
            names = sorted((n for n in z.namelist()
                            if re.match(r"ppt/slides/slide\d+\.xml$", n)),
                           key=lambda n: int(re.search(r"(\d+)", n).group(1)))
            for i, name in enumerate(names, 1):
                raw = z.read(name).decode("utf-8", errors="replace")
                runs = re.findall(r"<a:t>(.*?)</a:t>", raw, flags=re.S)
                slides.append((f"slide {i}", norm_ws(html.unescape(" ".join(runs)))))
        return slides, "zip-xml"
    except Exception:
        return None, None


def extract_epub(path: Path):
    try:
        with zipfile.ZipFile(path) as z:
            docs = [n for n in z.namelist()
                    if n.lower().endswith((".xhtml", ".html", ".htm")) and "toc" not in n.lower()]
            # try spine order from the OPF
            opf = next((n for n in z.namelist() if n.lower().endswith(".opf")), None)
            if opf:
                raw = z.read(opf).decode("utf-8", errors="replace")
                hrefs = re.findall(r'href="([^"]+\.x?html?)"', raw)
                base = os.path.dirname(opf)
                ordered = [os.path.normpath(os.path.join(base, h)) for h in hrefs]
                ordered = [h for h in ordered if h in docs]
                docs = ordered + [d for d in docs if d not in ordered]
            units = []
            for i, name in enumerate(docs, 1):
                text = strip_html(z.read(name).decode("utf-8", errors="replace"))
                if text:
                    units.append((f"ch.{i}", text))
            return units, "epub-zip"
    except Exception:
        return None, None


# ---------------------------------------------------------- chunking

_SENT_SPLIT = re.compile(r"(?<=[.!?])\s+(?=[A-Z\u00c0-\u024f\"'\u201c(])")


def split_paragraphs(unit_label: str, text: str):
    for para in re.split(r"\n\s*\n", text):
        para = para.strip()
        if not para:
            continue
        words = para.split()
        if len(words) <= CHUNK_MAX_WORDS:
            yield unit_label, para
        else:  # very long paragraph: split on sentences
            sents, cur = _SENT_SPLIT.split(para), []
            n = 0
            for s in sents:
                cur.append(s)
                n += len(s.split())
                if n >= CHUNK_TARGET_WORDS:
                    yield unit_label, " ".join(cur)
                    cur, n = [], 0
            if cur:
                yield unit_label, " ".join(cur)


def chunk_units(units):
    """Pack (unit_label, paragraph) stream into chunks near CHUNK_TARGET_WORDS.
    Returns list of dicts: {loc, text, words}. loc spans pages, e.g. 'p.3-p.5'."""
    chunks, buf, buf_words, loc_start, loc_end = [], [], 0, None, None

    def flush():
        nonlocal buf, buf_words, loc_start, loc_end
        if not buf:
            return
        loc = loc_start if (loc_start == loc_end or not loc_end) else f"{loc_start}-{loc_end}"
        chunks.append({"loc": loc or "", "text": "\n\n".join(buf), "words": buf_words})
        buf, buf_words, loc_start, loc_end = [], 0, None, None

    for unit_label, text in units:
        for label, para in split_paragraphs(unit_label, text):
            pwords = len(para.split())
            if buf_words + pwords > CHUNK_MAX_WORDS and buf_words >= CHUNK_MIN_WORDS:
                flush()
            if not buf:
                loc_start = label
            buf.append(para)
            loc_end = label
            buf_words += pwords
            if buf_words >= CHUNK_TARGET_WORDS:
                flush()
    flush()
    # merge a tiny trailing chunk into its predecessor
    if len(chunks) >= 2 and chunks[-1]["words"] < CHUNK_MIN_WORDS:
        last = chunks.pop()
        chunks[-1]["text"] += "\n\n" + last["text"]
        chunks[-1]["words"] += last["words"]
        if last["loc"] and last["loc"] != chunks[-1]["loc"]:
            start = chunks[-1]["loc"].split("-")[0]
            end = last["loc"].split("-")[-1]
            chunks[-1]["loc"] = start if start == end else f"{start}-{end}"
    return chunks


# ---------------------------------------------------------- index build

def classify(path: Path):
    ext = path.suffix.lower()
    if ext in PDF_EXTS: return "pdf"
    if ext in DOCX_EXTS: return "docx"
    if ext in PPTX_EXTS: return "pptx"
    if ext in EPUB_EXTS: return "epub"
    if ext in HTML_EXTS: return "html"
    if ext in CSV_EXTS: return "csv"
    if ext in JSON_EXTS: return "json"
    if ext in TEXT_EXTS: return "text"
    if ext in AUDIO_EXTS: return "audio"
    if ext in VIDEO_EXTS: return "video"
    if ext in IMAGE_EXTS: return "image"
    if ext in SKIP_EXTS: return "skip"
    return "unknown"


def extract_source(path: Path, kind: str):
    """Return (units, extractor_name, status)."""
    try:
        if kind == "text":
            return extract_plain(path), "plain", "indexed"
        if kind == "html":
            return extract_html_file(path), "html", "indexed"
        if kind == "csv":
            return extract_csv_file(path), "csv", "indexed"
        if kind == "json":
            return extract_json_file(path), "json", "indexed"
        if kind == "pdf":
            units, name = extract_pdf(path)
            if units:
                return units, name, "indexed"
            return None, None, "needs_extraction"
        if kind == "docx":
            units, name = extract_docx(path)
            return (units, name, "indexed") if units else (None, None, "needs_extraction")
        if kind == "pptx":
            units, name = extract_pptx(path)
            return (units, name, "indexed") if units else (None, None, "needs_extraction")
        if kind == "epub":
            units, name = extract_epub(path)
            return (units, name, "indexed") if units else (None, None, "needs_extraction")
        if kind in ("audio", "video"):
            return None, None, "needs_transcription"
        if kind == "image":
            return None, None, "needs_extraction"
        if kind == "unknown":
            # try as UTF-8 text; if it decodes cleanly enough, index it
            raw = path.read_bytes()[:8192]
            try:
                raw.decode("utf-8")
                return extract_plain(path), "plain", "indexed"
            except UnicodeDecodeError:
                return None, None, "skipped"
        return None, None, "skipped"
    except Exception as e:
        return None, None, f"error: {type(e).__name__}: {e}"


def discover_files(folder: Path):
    out = []
    for root, dirs, files in os.walk(folder):
        dirs[:] = [d for d in dirs if not d.startswith(".") and d not in
                   {NOTEBOOK_DIR.strip("."), "node_modules", "__pycache__", ".git"}]
        if NOTEBOOK_DIR in Path(root).parts:
            continue
        for fn in sorted(files):
            if fn.startswith("."):
                continue
            p = Path(root) / fn
            try:
                if p.stat().st_size == 0 or p.stat().st_size > MAX_FILE_BYTES:
                    continue
            except OSError:
                continue
            out.append(p)
    return sorted(out)


def parse_marked_text(raw: str):
    """Split agent-supplied text on [[page N]] / [[slide N]] / [[ch N]] markers."""
    parts = re.split(r"\[\[\s*(page|slide|ch|chapter)\s+([^\]\s]+)\s*\]\]", raw)
    if len(parts) == 1:
        return [("", norm_ws(raw))]
    units = []
    if parts[0].strip():
        units.append(("", norm_ws(parts[0])))
    label_map = {"page": "p.", "slide": "slide ", "ch": "ch.", "chapter": "ch."}
    for i in range(1, len(parts), 3):
        kind, num, body = parts[i], parts[i + 1], parts[i + 2]
        units.append((f"{label_map[kind]}{num}", norm_ws(body)))
    return [u for u in units if u[1]]


def load_registry(nb: Path):
    f = nb / "sources.json"
    if f.exists():
        return json.loads(f.read_text(encoding="utf-8"))
    return {"sources": []}


def save_all(nb: Path, registry, all_chunks):
    nb.mkdir(parents=True, exist_ok=True)
    (nb / "studio").mkdir(exist_ok=True)
    (nb / "notes").mkdir(exist_ok=True)
    (nb / "sources.json").write_text(
        json.dumps(registry, indent=1, ensure_ascii=False), encoding="utf-8")
    with open(nb / "chunks.jsonl", "w", encoding="utf-8") as f:
        for c in all_chunks:
            f.write(json.dumps(c, ensure_ascii=False) + "\n")
    meta = {"version": INDEX_VERSION,
            "built": datetime.now(timezone.utc).isoformat(),
            "chunks": len(all_chunks),
            "sources": len(registry["sources"])}
    (nb / "meta.json").write_text(json.dumps(meta, indent=1), encoding="utf-8")
    # invalidate search cache
    for c in (nb / "cache").glob("*.pkl") if (nb / "cache").exists() else []:
        c.unlink(missing_ok=True)


def next_source_id(registry):
    used = {int(s["id"][1:]) for s in registry["sources"] if re.match(r"S\d+$", s["id"])}
    n = 1
    while n in used:
        n += 1
    return f"S{n}"


def rebuild_chunks_for(source, units):
    chunks = chunk_units(units)
    out = []
    for i, c in enumerate(chunks):
        out.append({"id": f"{source['id']}#{i:03d}", "source_id": source["id"],
                    "source": source["title"], "loc": c["loc"],
                    "words": c["words"], "text": c["text"]})
    source["chunks"] = len(out)
    source["words"] = sum(c["words"] for c in out)
    return out


def cmd_build(folder: Path, from_text=None, ocr=False, ocr_lang="eng",
              ocr_dpi=300, ocr_engine="auto"):
    nb = folder / NOTEBOOK_DIR
    registry = load_registry(nb)
    by_path = {s["path"]: s for s in registry["sources"]}
    try:
        old_ver = json.loads((nb / "meta.json").read_text()).get("version", 0)
    except Exception:
        old_ver = 0
    force_reextract = old_ver != INDEX_VERSION
    if force_reextract and registry["sources"]:
        say(f"[index format v{old_ver} -> v{INDEX_VERSION}] re-extracting "
              "script-indexed sources (agent-supplied text is kept)")

    # keep chunks of sources we don't touch this run
    old_chunks = {}
    cj = nb / "chunks.jsonl"
    if cj.exists():
        with open(cj, encoding="utf-8") as f:
            for line in f:
                try:
                    c = json.loads(line)
                    old_chunks.setdefault(c["source_id"], []).append(c)
                except json.JSONDecodeError:
                    continue

    touched, all_chunks = set(), []

    if from_text:
        rel, textfile = from_text
        raw = Path(textfile).read_text(encoding="utf-8", errors="replace")
        units = parse_marked_text(raw)
        src = by_path.get(rel)
        p = folder / rel
        if not src:
            src = {"id": next_source_id(registry), "path": rel,
                   "title": title_from_path(rel), "type": classify(p) if p.exists() else "text",
                   "added": datetime.now(timezone.utc).isoformat()}
            registry["sources"].append(src)
            by_path[rel] = src
        src["sha256"] = sha256_file(p) if p.exists() else hashlib.sha256(raw.encode()).hexdigest()
        src["extractor"] = "agent"
        src["status"] = "indexed"
        src["units"] = len(units)
        all_chunks.extend(rebuild_chunks_for(src, units))
        touched.add(src["id"])
        print(f"[from-text] {src['id']} {rel}: {src['chunks']} chunks, {src['words']} words")
    else:
        files = discover_files(folder)
        present = set()
        for p in files:
            rel = str(p.relative_to(folder))
            present.add(rel)
            kind = classify(p)
            digest = sha256_file(p)
            src = by_path.get(rel)
            if (src and src.get("sha256") == digest
                    and src.get("status") == "indexed"
                    and not (force_reextract and src.get("extractor") != "agent")):
                continue  # unchanged
            if not src:
                src = {"id": next_source_id(registry), "path": rel,
                       "title": title_from_path(rel), "type": kind,
                       "added": datetime.now(timezone.utc).isoformat()}
                registry["sources"].append(src)
                by_path[rel] = src
            if src.get("extractor") == "agent" and src.get("sha256") == digest:
                continue  # agent-supplied text still valid
            src.update({"sha256": digest, "type": kind,
                        "bytes": p.stat().st_size})
            units, extractor, status = extract_source(p, kind)

            # Scanned-PDF detection: a text layer that thin is useless for Q&A.
            if kind == "pdf" and units:
                total_w = sum(len(t.split()) for _, t in units)
                if units and total_w / max(1, len(units)) < SCAN_WORDS_PER_PAGE:
                    src["note"] = (f"suspected scan: avg "
                                   f"{total_w / max(1, len(units)):.0f} words/page")
                    units, extractor, status = None, None, "needs_extraction"

            # OCR pathway for scanned PDFs and images.
            if status == "needs_extraction" and kind in ("pdf", "image") and ocr:
                if run_ocr:
                    try:
                        pages = run_ocr(p, lang=ocr_lang, dpi=ocr_dpi,
                                        engine=ocr_engine)
                        units = [(f"p.{n}", t) for n, t in pages if t.strip()]
                        if units:
                            extractor, status = "ocr", "indexed"
                            src["note"] = (f"OCR'd (engine={ocr_engine}, "
                                           f"{ocr_lang}, {ocr_dpi}dpi)")
                        else:
                            src["note"] = "OCR produced no text"
                    except Exception as e:
                        src["note"] = f"OCR failed: {e}"
                else:
                    src["note"] = ("OCR unavailable: ocr_pdf.py missing next "
                                   "to build_index.py")

            src["status"] = status
            src["extractor"] = extractor
            touched.add(src["id"])
            if units:
                src["units"] = len(units)
                all_chunks.extend(rebuild_chunks_for(src, units))
                say(f"[indexed] {src['id']} {rel} ({extractor}): "
                      f"{src['chunks']} chunks, {src['words']} words")
            else:
                src["chunks"], src["words"] = 0, 0
                say(f"[{status}] {src['id']} {rel}")
        # drop sources whose files vanished (but keep agent-added virtual sources)
        registry["sources"] = [s for s in registry["sources"]
                               if s["path"] in present or s.get("extractor") == "agent"]

    # carry over untouched sources' chunks
    kept_ids = {s["id"] for s in registry["sources"]}
    for sid, chunks in old_chunks.items():
        if sid in kept_ids and sid not in touched:
            all_chunks.extend(chunks)
    all_chunks.sort(key=lambda c: (int(c["source_id"][1:]), c["id"]))

    save_all(nb, registry, all_chunks)
    print_status(folder)


def print_status(folder: Path):
    nb = folder / NOTEBOOK_DIR
    if not (nb / "sources.json").exists():
        say("No index yet. Run: python build_index.py <folder>")
        return
    registry = load_registry(nb)
    total_words = sum(s.get("words", 0) for s in registry["sources"])
    total_chunks = sum(s.get("chunks", 0) for s in registry["sources"])
    say(f"\n=== Notebook index: {folder} ===")
    say(f"{len(registry['sources'])} sources · {total_chunks} chunks · {total_words:,} words\n")
    for s in sorted(registry["sources"], key=lambda x: int(x["id"][1:])):
        flag = "" if s.get("status") == "indexed" else f"  <-- {s.get('status')}"
        say(f"  {s['id']:<4} {s.get('status','?'):<18} {s.get('words',0):>9,}w  "
              f"{s['path']}{flag}")
    pending = [s for s in registry["sources"] if s.get("status") != "indexed"]
    if pending:
        say("\nSources needing attention:")
        for s in pending:
            note = f" — {s['note']}" if s.get("note") else ""
            say(f"  - {s['path']} [{s.get('status')}]{note}")
        say("\nRemedies: (a) re-run with --ocr to OCR scanned PDFs/images "
              "(needs tesseract; add --ocr-lang for non-English), or "
              "(b) extract/transcribe the text yourself and run "
              "--from-text <relpath> <extracted.txt>.")


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("folder", help="Folder containing the sources")
    ap.add_argument("--status", action="store_true", help="Show index status only")
    ap.add_argument("--from-text", nargs=2, metavar=("RELPATH", "TEXTFILE"),
                    help="Register agent-extracted text for a source")
    ap.add_argument("--ocr", action="store_true",
                    help="OCR scanned PDFs and images that lack a usable text layer")
    ap.add_argument("--ocr-lang", default="eng",
                    help="tesseract language(s) for --ocr, e.g. eng, deu, eng+bul")
    ap.add_argument("--ocr-dpi", type=int, default=300,
                    help="Rasterization DPI for --ocr (default 300)")
    ap.add_argument("--ocr-engine", default="auto",
                    choices=["auto", "tesseract", "easyocr"],
                    help="OCR engine for --ocr (auto = easyocr when installed)")
    args = ap.parse_args()
    folder = Path(args.folder).resolve()
    if not folder.is_dir():
        sys.exit(f"Folder not found: {folder}\n"
                 "If the path contains spaces, wrap it in quotes, e.g.\n"
                 f'  python {Path(__file__).name} "C:\\path with spaces\\books" ...')
    if args.status:
        print_status(folder)
    else:
        cmd_build(folder, from_text=args.from_text, ocr=args.ocr,
                  ocr_lang=args.ocr_lang, ocr_dpi=args.ocr_dpi,
                  ocr_engine=args.ocr_engine)


if __name__ == "__main__":
    try:
        main()
    except BrokenPipeError:
        try:
            sys.stdout.close()
        except Exception:
            pass
